#!/usr/bin/env python3
"""
Convert the RiskRunners presentation markdown to PowerPoint.
Requires: pip install python-pptx Pillow
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
import re

# Colors from the markdown theme
DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
RED = RGBColor(0xCC, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def add_title_slide(prs, title, subtitle=None):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(32)
        p.font.color.rgb = RED
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines, image_path=None):
    """Add a content slide with optional image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Determine content width based on image presence
    content_width = Inches(5.5) if image_path else Inches(9)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), content_width, Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    first = True
    for line in content_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        
        # Handle bullet points
        if line.startswith('* ') or line.startswith('- '):
            p.text = '• ' + line[2:]
            p.level = 0
        elif line.startswith('  * ') or line.startswith('  - '):
            p.text = '  ◦ ' + line[4:]
            p.level = 1
        elif line.startswith('    * ') or line.startswith('    - '):
            p.text = '    ▪ ' + line[6:]
            p.level = 2
        elif re.match(r'^\d+\.\s+', line):
            p.text = line
            p.level = 0
        else:
            p.text = line
        
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        
        # Bold text between ** **
        if '**' in p.text:
            # Simple approach - just remove markers for now
            p.text = p.text.replace('**', '')
    
    # Add image if provided
    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, Inches(6.2), Inches(1.5), width=Inches(3.5))
        except Exception as e:
            print(f"Could not add image {image_path}: {e}")
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "What Was - Is For What Will Be"
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.8), Inches(9), Inches(4)).table
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header.replace('**', '')
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xAB, 0xBA, 0xC9)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text.replace('**', '')
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
    
    return slide

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    base_path = os.path.dirname(os.path.abspath(__file__))
    
    # Slide 1: Title
    slide = add_title_slide(prs, "The Big Trust", "What Was - Is For What Will Be")
    # Add additional info
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "A Founder's Journey: From Risk Runners to Risk Management"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Presented to the Risk Runners Club, University of Arizona"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "February 17, 2026"
    p.font.size = Pt(18)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: The Elephant in the Room
    add_content_slide(prs, "The Elephant in the Room", [
        "I founded this club, but I never passed an actuarial exam.",
        "",
        "• The Reality: I oriented my entire life toward actuarial science,",
        "  but the path wasn't linear.",
        "• The Outcome: Orienting toward this discipline changed my life,",
        "  even without the credentials.",
        "• The Goal Today: To share the \"Big Trust\"—how understanding",
        "  the architecture of risk can build a career, whether you pass",
        "  the exams or not."
    ], os.path.join(base_path, "the-career-pivot.png"))
    
    # Slide 3: Origins
    add_content_slide(prs, "Origins: 10 Fingers vs. 100,000", [
        "Before the math, there was Language & Code.",
        "",
        "• 2015: Wrote a book on AI Philosophy.",
        "• The Tech: Created LinguaLint (Natural Language Processing).",
        "  - Scraped 1,500 news articles/day.",
        "  - Extracted event codes to predict stock market trends.",
        "• The Realization: It was a \"lightweight LLM middleware.\"",
        "  - My 10 fingers coding from scratch vs. the 100,000 fingers",
        "    that built ChatGPT.",
        "• The Shift: Moved from UC Davis (Ag) → Poetry → U of A",
        "  (Math/Finance) to find something \"eternal\" amidst rapid",
        "  tech changes."
    ], os.path.join(base_path, "hands.png"))
    
    # Slide 4: Educational Crossroads
    add_content_slide(prs, "The Educational Crossroads", [
        "Why U of A?",
        "",
        "• The Choice: ASU (Established Actuarial Program) vs.",
        "  U of A (College Town Culture).",
        "• The Gap: ASU had the funding, the WSIA events, and the",
        "  Casualty Actuarial Society (CAS) connection. U of A did not.",
        "• The Solution: Founded Risk Runners.",
        "  - Bridging the gap.",
        "  - Driving students to WSIA events.",
        "  - Creating a community from scratch (0 to 1)."
    ])
    
    # Slide 5: The Overload
    add_content_slide(prs, "The Overload & The Hubris", [
        "I tried to do it all at once:",
        "",
        "1. Math Degree (PDEs, Stochastic Processes).",
        "2. Finance Program (Eller).",
        "3. Running a Tech Company (Sensutec).",
        "4. Latin & Systems Engineering.",
        "5. Running the Risk Runners Club.",
        "6. Exercise",
        "7. Night Life",
        "",
        "The Result: I took the Financial Math exam to \"see how hard",
        "it was.\" I failed.",
        "",
        "But in failing the exam, I got to see the bigger picture."
    ])
    
    # Slide 6: Nautilus Internship
    add_content_slide(prs, "The Actuarial Internship: Nautilus Insurance", [
        "Thanks to Brent Carr @ Nautilus / Connection via WSIA",
        "",
        "I applied my NLP software to the insurance world.",
        "",
        "• The Project: Scraping Yelp reviews for hotel/motels.",
        "• The Insight: \"Great mimosa!\" + No alcohol indemnification",
        "  = Risk Gap.",
        "• The Win-Win-Win:",
        "  1. Underwriter: Knows the risk before the lawyers do.",
        "  2. Insurer: Gets higher premiums for proper coverage.",
        "  3. Client: Is actually covered for their reality."
    ], os.path.join(base_path, "nautilus-nlp-project.png"))
    
    # Slide 7: Enterprise Architecture
    add_content_slide(prs, "The Enterprise Architecture", [
        "What I learned inside the building.",
        "",
        "Insurance isn't just math; it's a physical structure of talent.",
        "",
        "• Floor 1: Claims & Database Admins (The Foundation).",
        "• Floor 2: Underwriters & Actuaries (The Risk Assessment).",
        "• Floor 3: Financial Analysts & CEO (The Investment/Strategy).",
        "",
        "Key Takeaway: Insurance doesn't make money solely on",
        "premiums; it makes money by investing the float."
    ], os.path.join(base_path, "enterprise-arch.png"))
    
    # Slide 8: Tale of Two Societies (Table)
    add_table_slide(prs, "The Vision: A Tale of Two Societies",
        ["CAS (Casualty)", "SOA (Society of Actuaries)"],
        [
            ["Focus: Property/Product", "Focus: Health/Life"],
            ["Risk: Before/After Lifecycle", "Risk: During Lifecycle"],
            ["Analogy: The Lawyer", "Analogy: The Doctor"],
            ["School: ASU (State School)", "School: U of A (University)"]
        ]
    )
    # Add goal text to last slide
    goal_box = prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.8))
    tf = goal_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Goal: Give high schoolers a clear choice of path, just like choosing between Med School and Law School."
    p.font.size = Pt(16)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 9: Career Trajectory
    add_content_slide(prs, "Career Trajectory: The Pivot", [
        "Post-2020: Adapting to the World",
        "",
        "1. CPG Analytics (Remote Work)",
        "   • Started with simple tools.",
        "   • Mastered Tableau (Business Intelligence).",
        "   • Action: Taught others how to use Kaggle data + AI",
        "     to build portfolios.",
        "",
        "2. Veteran Affairs (Contractor)",
        "   • Patient Generated Health Data (PGHD).",
        "   • Wearables (Fitbit, Dexcom) predicting health risks.",
        "   • The Actuarial Link: Using continuous IoT data to",
        "     predict insurance risk (Life/Health)."
    ])
    
    # Slide 10: Captive Insurance
    add_content_slide(prs, "Current Frontier: Captive Insurance", [
        "Merging Housing & Risk",
        "",
        "My market research led me to the housing crisis and",
        "Prefab Manufacturing.",
        "",
        "• The Strategy: Applying the \"Insurance Building\" model",
        "  to a smaller scale.",
        "• Captive Insurance:",
        "  - Like 12 doctors self-insuring to pool risk.",
        "  - Capturing the premium + Investing the float.",
        "• The Focus: Captive Management for Prefab Manufacturers",
        "  (captive.integralmass.com)"
    ], os.path.join(base_path, "captive-iMASS-logo.png"))
    
    # Slide 11: What I Value
    add_content_slide(prs, "What I Value (vs. Advice)", [
        "Academia vs. Industry",
        "• Academia: Can trap you in a cycle of paying for credentials.",
        "• Industry: Pays you for your time and values work experience",
        "  and tool mastery (tech stack) over pure theory.",
        "",
        "Work-Life Balance",
        "• I value leaving work at work.",
        "• Time for: Hiking, Poetry, Pickleball, Networking.",
        "• We work to live, we don't live to work."
    ])
    
    # Slide 12: Resources
    add_content_slide(prs, "Resources & Community", [
        "Where to look next",
        "",
        "• AI Trailblazers:",
        "  - Aaron Eden & Maria Eden.",
        "  - Bridging the gap between AI-literate and non-literate.",
        "  - Training mentors to train apprentices.",
        "",
        "• Casualty Actuaries of the Desert States.",
        "• Arizona Captive Insurance Association (AZCIA).",
        "• Venture Café in Phoenix.",
        "",
        "Networking is critical."
    ])
    
    # Slide 13: Final Thought
    slide = add_title_slide(prs, "The Final Thought")
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(4))
    tf = content_box.text_frame
    p = tf.paragraphs[0]
    p.text = "\"What was - is for what will be.\""
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "I may go back and take those first two exams because they open doors."
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "But my journey taught me that Risk Management is broader than just a test score."
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "It's about: Networking • Technology • Trust (In your path)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank you"
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Jefferson Richards | 520.981.3639 | jefferson@richards.plus"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "jefferson.cloud | richards.systems | richards.plus"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER
    
    # Save
    output_path = os.path.join(base_path, "2026-01-27_RiskRunners-Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    main()
